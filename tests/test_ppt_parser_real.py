from pathlib import Path

from pptx import Presentation

from lodestar_veritas.parsers.ppt_parser import PPTParser


def test_ppt_parser_extracts_slide_text(tmp_path: Path):
    ppt_path = tmp_path / "sample.pptx"

    presentation = Presentation()
    slide = presentation.slides.add_slide(presentation.slide_layouts[0])
    slide.shapes.title.text = "Financial Summary"
    slide.placeholders[1].text = "Revenue increased by 25% in 2024."

    presentation.save(ppt_path)

    parser = PPTParser()
    result = parser.parse(str(ppt_path))

    assert result["metadata"]["parser"] == "PPTParser"
    assert result["metadata"]["slide_count"] == 1
    assert "Revenue increased" in result["content"]