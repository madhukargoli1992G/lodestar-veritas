from pathlib import Path

from pptx import Presentation


class PPTParser:
    """
    Real PPTX parser.

    Extracts slide text and slide metadata.
    """

    def parse(self, file_path: str) -> dict:
        path = Path(file_path)
        presentation = Presentation(file_path)

        slides = []
        text_parts = []

        for slide_index, slide in enumerate(presentation.slides, start=1):
            slide_text_parts = []

            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    text = shape.text.strip()
                    if text:
                        slide_text_parts.append(text)

            slide_text = "\n".join(slide_text_parts).strip()

            slides.append(
                {
                    "slide_number": slide_index,
                    "text": slide_text,
                }
            )

            if slide_text:
                text_parts.append(f"[Slide {slide_index}]\n{slide_text}")

        content = "\n\n".join(text_parts).strip()

        return {
            "content": content,
            "slides": slides,
            "metadata": {
                "source": file_path,
                "file_name": path.name,
                "parser": "PPTParser",
                "slide_count": len(slides),
            },
        }